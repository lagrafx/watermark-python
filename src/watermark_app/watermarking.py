"""Watermark utilities for supported document types."""

from __future__ import annotations

import hashlib
import io
import tempfile
from pathlib import Path

SUPPORTED_EXTENSIONS = {".docx", ".docm", ".xlsx", ".xlsm", ".pptx", ".pptm", ".pdf"}
WATERMARK_OPACITY = 0.35
WATERMARK_MARKER = "watermark-python|behind-text-v2"
EMU_PER_INCH = 914400
LEGACY_WORD_WATERMARK_WIDTH_EMU = 6 * EMU_PER_INCH
WORD_WATERMARK_WIDTH_INCHES = 6.5


def is_supported_extension(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def apply_watermark(source_path: Path, output_path: Path, watermark_png_path: Path) -> None:
    with tempfile.TemporaryDirectory(prefix="watermark_asset_") as tmp_dir:
        effective_watermark = Path(tmp_dir) / "watermark.png"
        _write_opacity_adjusted_png(watermark_png_path, effective_watermark, WATERMARK_OPACITY)

        ext = source_path.suffix.lower()
        if ext in {".docx", ".docm"}:
            _watermark_word(source_path, output_path, effective_watermark, watermark_png_path)
            return
        if ext in {".xlsx", ".xlsm"}:
            _watermark_excel(source_path, output_path, effective_watermark)
            return
        if ext in {".pptx", ".pptm"}:
            _watermark_powerpoint(source_path, output_path, effective_watermark)
            return
        if ext == ".pdf":
            _watermark_pdf(source_path, output_path, effective_watermark)
            return
        raise ValueError(f"Unsupported file extension: {ext}")


def _write_opacity_adjusted_png(source_path: Path, output_path: Path, opacity: float) -> None:
    from PIL import Image

    opacity = max(0.0, min(1.0, opacity))
    image = Image.open(source_path).convert("RGBA")
    alpha = image.getchannel("A").point(lambda value: int(value * opacity))
    image.putalpha(alpha)
    image.save(output_path, "PNG")


def _sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _watermark_word(
    source_path: Path,
    output_path: Path,
    watermark_png_path: Path,
    legacy_watermark_png_path: Path,
) -> None:
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches

    document = Document(str(source_path))
    removable_image_hashes = {
        _sha256_bytes(legacy_watermark_png_path.read_bytes()),
        _sha256_bytes(watermark_png_path.read_bytes()),
    }

    def add_header_watermark(header) -> None:  # noqa: ANN001
        _remove_existing_header_images(header, removable_image_hashes)
        paragraph = header.add_paragraph()
        run = paragraph.add_run()
        run.add_picture(str(watermark_png_path), width=Inches(WORD_WATERMARK_WIDTH_INCHES))
        drawing = run._r.xpath("./w:drawing")[0]  # noqa: SLF001
        inline = drawing.xpath("./wp:inline")[0]
        doc_pr = inline.xpath("./wp:docPr")[0]
        doc_pr.set("name", WATERMARK_MARKER)
        doc_pr.set("descr", WATERMARK_MARKER)
        inline.tag = qn("wp:anchor")
        inline.set("distT", "0")
        inline.set("distB", "0")
        inline.set("distL", "0")
        inline.set("distR", "0")
        inline.set("simplePos", "0")
        inline.set("relativeHeight", "251659264")
        inline.set("behindDoc", "1")
        inline.set("locked", "0")
        inline.set("layoutInCell", "1")
        inline.set("allowOverlap", "1")

        simple_pos = OxmlElement("wp:simplePos")
        simple_pos.set("x", "0")
        simple_pos.set("y", "0")

        position_h = OxmlElement("wp:positionH")
        position_h.set("relativeFrom", "page")
        align_h = OxmlElement("wp:align")
        align_h.text = "center"
        position_h.append(align_h)

        position_v = OxmlElement("wp:positionV")
        position_v.set("relativeFrom", "page")
        align_v = OxmlElement("wp:align")
        align_v.text = "center"
        position_v.append(align_v)

        wrap_none = OxmlElement("wp:wrapNone")
        inline.insert(0, simple_pos)
        inline.insert(1, position_h)
        inline.insert(2, position_v)
        inline.insert(3, wrap_none)

    for section in document.sections:
        add_header_watermark(section.header)
        if getattr(section, "different_first_page_header_footer", False):
            add_header_watermark(section.first_page_header)
        if getattr(document.settings, "odd_and_even_pages_header_footer", False):
            add_header_watermark(section.even_page_header)
    document.save(str(output_path))


def _remove_existing_header_images(
    header,  # noqa: ANN001
    removable_image_hashes: set[str],
) -> None:
    """Remove tagged watermarks and proven legacy watermark images from Word headers."""
    for paragraph in list(header.paragraphs):
        for run in list(paragraph.runs):
            if _is_tagged_watermark_run(run) or _is_matching_legacy_watermark_run(
                header,
                run,
                removable_image_hashes,
            ):
                run._element.getparent().remove(run._element)  # noqa: SLF001
        if not paragraph.text and not paragraph._p.xpath(".//w:drawing | .//w:pict"):  # noqa: SLF001
            parent = paragraph._p.getparent()  # noqa: SLF001
            if parent is not None and len(header.paragraphs) > 1:
                parent.remove(paragraph._p)  # noqa: SLF001


def _is_tagged_watermark_run(run) -> bool:  # noqa: ANN001
    for doc_pr in run._r.xpath(".//wp:docPr | .//pic:cNvPr"):  # noqa: SLF001
        values = [doc_pr.get("name", ""), doc_pr.get("descr", ""), doc_pr.get("title", "")]
        if any(WATERMARK_MARKER in value for value in values):
            return True
    return False


def _is_matching_legacy_watermark_run(
    header,  # noqa: ANN001
    run,  # noqa: ANN001
    removable_image_hashes: set[str],
) -> bool:
    if not _has_legacy_watermark_size(run):
        return False

    for relationship_id in run._r.xpath(".//a:blip/@r:embed"):  # noqa: SLF001
        image_part = header.part.related_parts.get(relationship_id)
        if image_part is None:
            continue
        if _sha256_bytes(image_part.blob) in removable_image_hashes:
            return True
    return False


def _has_legacy_watermark_size(run) -> bool:  # noqa: ANN001
    tolerance = int(0.25 * EMU_PER_INCH)
    min_width = LEGACY_WORD_WATERMARK_WIDTH_EMU - tolerance
    max_width = LEGACY_WORD_WATERMARK_WIDTH_EMU + tolerance
    for extent in run._r.xpath(".//wp:extent"):  # noqa: SLF001
        try:
            width = int(extent.get("cx", "0"))
        except ValueError:
            continue
        if min_width <= width <= max_width:
            return True
    return False


def _watermark_excel(source_path: Path, output_path: Path, watermark_png_path: Path) -> None:
    from openpyxl import load_workbook
    from openpyxl.drawing.image import Image as XLImage

    keep_vba = source_path.suffix.lower() == ".xlsm"
    workbook = load_workbook(filename=str(source_path), keep_vba=keep_vba)
    for sheet in workbook.worksheets:
        image = XLImage(str(watermark_png_path))
        sheet.add_image(image, "A1")
    workbook.save(str(output_path))


def _watermark_powerpoint(source_path: Path, output_path: Path, watermark_png_path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation(str(source_path))
    for slide in presentation.slides:
        # Scale watermark to ~60% of slide width and center it.
        target_width = int(presentation.slide_width * 0.6)
        picture = slide.shapes.add_picture(
            str(watermark_png_path),
            left=0,
            top=0,
            width=target_width,
        )
        picture.left = int((presentation.slide_width - picture.width) / 2)
        picture.top = int((presentation.slide_height - picture.height) / 2)
        picture.name = WATERMARK_MARKER
        for c_nv_pr in picture._element.xpath(".//p:cNvPr"):  # noqa: SLF001
            c_nv_pr.set("descr", WATERMARK_MARKER)
        shape_tree = slide.shapes._spTree  # noqa: SLF001
        shape_tree.remove(picture._element)  # noqa: SLF001
        shape_tree.insert(2, picture._element)  # noqa: SLF001
    presentation.save(str(output_path))


def _watermark_pdf(source_path: Path, output_path: Path, watermark_png_path: Path) -> None:
    from pypdf import PdfReader, PdfWriter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    def _overlay_page_bytes(page_width: float, page_height: float) -> bytes:
        buffer = io.BytesIO()
        pdf = canvas.Canvas(buffer, pagesize=(page_width, page_height))
        img = ImageReader(str(watermark_png_path))
        img_width, img_height = img.getSize()

        # Fit watermark to max 70% of page width/height while preserving aspect ratio.
        max_width = page_width * 0.7
        max_height = page_height * 0.7
        scale = min(max_width / img_width, max_height / img_height)
        draw_width = img_width * scale
        draw_height = img_height * scale
        x = (page_width - draw_width) / 2
        y = (page_height - draw_height) / 2

        pdf.drawImage(
            str(watermark_png_path),
            x,
            y,
            width=draw_width,
            height=draw_height,
            mask="auto",
            preserveAspectRatio=True,
        )
        pdf.showPage()
        pdf.save()
        buffer.seek(0)
        return buffer.read()

    reader = PdfReader(str(source_path))
    writer = PdfWriter()
    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        overlay_reader = PdfReader(io.BytesIO(_overlay_page_bytes(width, height)))
        watermark_page = overlay_reader.pages[0]
        watermark_page.merge_page(page)
        writer.add_page(watermark_page)

    with output_path.open("wb") as out_file:
        writer.write(out_file)
