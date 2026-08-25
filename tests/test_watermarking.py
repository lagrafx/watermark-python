import zipfile
from pathlib import Path

from PIL import Image

from watermark_app.watermarking import apply_watermark


def _make_watermark(path: Path) -> None:
    image = Image.new("RGBA", (400, 120), (255, 0, 0, 255))
    image.save(path)


def test_word_watermark_is_inserted_as_first_page_vml_watermark(tmp_path: Path) -> None:
    from docx import Document

    source = tmp_path / "source.docx"
    output = tmp_path / "output.docx"
    watermark = tmp_path / "watermark.png"
    _make_watermark(watermark)

    document = Document()
    document.add_paragraph("Body text")
    document.save(source)

    apply_watermark(source, output, watermark)
    Document(str(output))

    with zipfile.ZipFile(output) as package:
        headers = [
            package.read(name).decode("utf-8")
            for name in package.namelist()
            if name.startswith("word/header") and name.endswith(".xml")
        ]

    assert headers
    assert any("<v:shape" in header for header in headers)
    assert sum(header.count("WatermarkPythonFirstPageBehindText") for header in headers) == 1
    assert any("watermark-python|first-page-behind-text-v3" in header for header in headers)


def test_word_repair_removes_matching_legacy_watermark_but_keeps_header_logo(
    tmp_path: Path,
) -> None:
    from docx import Document
    from docx.shared import Inches

    source = tmp_path / "source.docx"
    output = tmp_path / "output.docx"
    logo = tmp_path / "logo.png"
    watermark = tmp_path / "watermark.png"
    Image.new("RGBA", (80, 40), (0, 0, 255, 255)).save(logo)
    _make_watermark(watermark)

    document = Document()
    document.add_paragraph("Body text")
    header = document.sections[0].header
    logo_run = header.paragraphs[0].add_run()
    logo_run.add_picture(str(logo), width=Inches(1.0))
    old_watermark_run = header.paragraphs[0].add_run()
    old_watermark_run.add_picture(str(watermark), width=Inches(6.0))
    document.save(source)

    apply_watermark(source, output, watermark)
    Document(str(output))

    with zipfile.ZipFile(output) as package:
        header_xml = "\n".join(
            package.read(name).decode("utf-8")
            for name in package.namelist()
            if name.startswith("word/header") and name.endswith(".xml")
        )

    assert header_xml.count("<wp:inline") == 1
    assert header_xml.count("WatermarkPythonFirstPageBehindText") == 1
    assert "watermark-python|first-page-behind-text-v3" in header_xml


def test_powerpoint_watermark_is_first_slide_shape_after_group_properties(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    source = tmp_path / "source.pptx"
    output = tmp_path / "output.pptx"
    watermark = tmp_path / "watermark.png"
    _make_watermark(watermark)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(100, 100, 300, 100).text = "Body text"
    presentation.save(source)

    apply_watermark(source, output, watermark)

    with zipfile.ZipFile(output) as package:
        slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")

    pic_index = slide_xml.index("<p:pic>")
    text_box_index = slide_xml.index("<p:sp>")
    assert pic_index < text_box_index
    assert "watermark-python|first-page-behind-text-v3" in slide_xml


def test_excel_embedded_watermark_png_is_semi_transparent(tmp_path: Path) -> None:
    from openpyxl import Workbook

    source = tmp_path / "source.xlsx"
    output = tmp_path / "output.xlsx"
    watermark = tmp_path / "watermark.png"
    _make_watermark(watermark)

    workbook = Workbook()
    workbook.active["A1"] = "Body text"
    workbook.save(source)

    apply_watermark(source, output, watermark)

    with zipfile.ZipFile(output) as package:
        media_names = [name for name in package.namelist() if name.startswith("xl/media/")]
        assert media_names
        with package.open(media_names[0]) as media_file:
            embedded = Image.open(media_file).convert("RGBA")
            alpha = embedded.getchannel("A").getextrema()

    assert alpha[1] < 255
